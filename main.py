import flet as ft
import json
import urllib.request
import threading
import time
import os
import uuid

# For native Windows picker (workaround for desktop)
def native_pick_files():
    try:
        import tkinter as tk
        from tkinter import filedialog
        root = tk.Tk()
        root.withdraw()
        # Ensure it stays on top
        root.attributes("-topmost", True)
        files = filedialog.askopenfilenames(
            title="Select Files for Bulk Upload",
            filetypes=[("Lyrics Documents", "*.txt;*.pdf;*.pptx"), ("All Files", "*.*")]
        )
        root.destroy()
        return files
    except:
        return []

# ===================== IN-MEMORY DATA STORE =====================
SONGS = []

FIRESTORE_URL = "https://gggm-admin.rajhanoch24.workers.dev"

def seed_default_songs():
    global SONGS
    if len(SONGS) > 0:
        return
    SONGS = [
        {
            "id": "ta_172",
            "title": "172. கர்த்தாவே தேவர்களில்",
            "language": "tamil",
            "is_favorite": False,
            "lyrics": (
                "கர்த்தாவே தேவர்களில் உமக்கொப்பனவர் யார்?\n"
                "வானத்திலும் பூமியிலும் உமக்கொப்பானவர் யார்?\n\n"
                "உமக்கொப்பனவர் யார்? உமக்கொப்பனவர் யார்?\n"
                "வானத்திலும் பூமியிலும் உமக்கொப்பானவர் யார்?\n\n"
                "1. செங்கடலை நீர் பிளந்து\n"
                "உந்தன் ஜனங்களை நடத்திச் சென்றீர்\n"
                "நீர் நல்லவர் சர்வ வல்லவர்\n"
                "என்றும் வாக்கு மாறாதவர் (2)\n\n"
                "2. தூதர்கள் உண்ணும் உணவால்\n"
                "உந்தன் ஜனங்களை போஷித்தீரே\n"
                "உம்மைப்போல யாருண்டு\n"
                "இந்த ஜனங்களை நேசித்திட (2)"
            ),
        },
        {
            "id": "te_1",
            "title": "అగ్ని మండించు - Agni Mandinchu",
            "language": "telugu",
            "is_favorite": False,
            "lyrics": (
                "అగ్ని మండించు – నాలో అగ్ని మండించు (2)\n"
                "పరిశుద్ధాత్ముడా – నాలో అగ్ని మండించు (2)\n\n"
                "Agni Mandinchu – Naalo Agni Mandinchu (2)\n"
                "Parishuddhaathmudaa – Naalo Agni Mandinchu (2)"
            ),
        },
    ]

def get_filtered(lang, query=""):
    result = []
    for s in SONGS:
        if s["language"] != lang:
            continue
        if query and query.lower() not in s["title"].lower():
            continue
        result.append(s)
    return sorted(result, key=lambda x: x["title"])

def cloud_sync():
    global SONGS
    count = 0
    try:
        # Fetch from our new Cloudflare /songs endpoint
        req = urllib.request.Request(f"{FIRESTORE_URL}/songs")
        with urllib.request.urlopen(req, timeout=15) as resp:
            body = resp.read().decode("utf-8")
            cloud_list = json.loads(body)
            
            # Simple list of dicts directly from Cloudflare D1
            if isinstance(cloud_list, list):
                fav_ids = {s["id"] for s in SONGS if s.get("is_favorite", False)}
                new_songs = []
                for s in cloud_list:
                    song = {
                        "id": str(s.get("id")),
                        "title": s.get("title", "Untitled"),
                        "language": s.get("language", "tamil").lower(),
                        "lyrics": s.get("lyrics", ""),
                        "is_favorite": str(s.get("id")) in fav_ids,
                    }
                    new_songs.append(song)
                
                if new_songs:
                    SONGS = new_songs
                    count = len(SONGS)
    except Exception as e:
        print(f"Sync error: {e}")
    return count

# ===================== MAIN APP =====================
def main(page: ft.Page):
    page.window_icon = os.path.join(os.getcwd(), "assets", "icon.png")
    page.title = "GGGM"
    page.theme_mode = ft.ThemeMode.LIGHT
    page.bgcolor = "#F5F7FA"

    # ---- File Picker Service ----
    picker = ft.FilePicker()
    
    async def process_bulk_upload(e):
        selected_paths = []
        
        # Check platform
        is_android = page.platform == ft.PagePlatform.ANDROID or page.platform == ft.PagePlatform.IOS
        
        if is_android:
            # Use standard Flet Picker on Mobile
            if picker not in page.overlay:
                page.overlay.append(picker)
                page.update()
            
            res = await picker.pick_files(
                allow_multiple=True, 
                allowed_extensions=["txt", "pdf", "pptx"],
                file_type=ft.FilePickerFileType.CUSTOM
            )
            if res:
                selected_paths = [f.path for f in res if f.path]
        else:
            # Use native Windows picker on Desktop to avoid "Unknown Control" error
            status_in = getattr(page, "admin_status", None)
            if status_in:
                status_in.value = "Opening Windows File Selector..."
                page.update()
            
            # This is a blocking call, so we run it in a way that doesn't freeze Flet UI
            selected_paths = native_pick_files()

        if not selected_paths:
            status_in = getattr(page, "admin_status", None)
            if status_in:
                status_in.value = ""
                page.update()
            return
        
        status_in = getattr(page, "admin_status", None)
        if status_in:
            status_in.value = f"Processing {len(selected_paths)} files..."
            status_in.color = "blue"
            page.update()

        def background_processing():
            success_count = 0
            for path in selected_paths:
                try:
                    name = os.path.basename(path)
                    title = os.path.splitext(name)[0]
                    content = ""
                    
                    if name.lower().endswith(".txt"):
                        with open(path, 'r', encoding='utf-8', errors='ignore') as file:
                            content = file.read()
                    elif name.lower().endswith(".pdf"):
                        try:
                            from pypdf import PdfReader
                            reader = PdfReader(path)
                            for page_pdf in reader.pages:
                                content += page_pdf.extract_text() + "\n"
                        except ImportError:
                            content = "[PDF Library Error: Please rebuild APK with requirements]"
                    elif name.lower().endswith(".pptx"):
                        try:
                            from pptx import Presentation
                            prs = Presentation(path)
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        content += shape.text + "\n"
                        except ImportError:
                            content = "[PPTX Library Error: Please rebuild APK with requirements]"
                    
                    if content.strip():
                        # Push to cloud
                        sid = f"bk_{uuid.uuid4().hex[:8]}"
                        document = {
                            "fields": {
                                "title": {"stringValue": title},
                                "language": {"stringValue": getattr(page, "admin_lang", "tamil")},
                                "number": {"stringValue": ""},
                                "lyrics": {"stringValue": content.replace("\n", "\\n")},
                                "category": {"stringValue": "General"},
                                "composer": {"stringValue": "Unknown"}
                            }
                        }
                        req = urllib.request.Request(
                            f"{FIRESTORE_URL}/{sid}",
                            data=json.dumps(document).encode("utf-8"),
                            headers={"Content-Type": "application/json"},
                            method="PATCH"
                        )
                        with urllib.request.urlopen(req) as r:
                            if r.status == 200:
                                success_count += 1
                except Exception as ex:
                    print(f"Error processing {path}: {ex}")
            
            if status_in:
                status_in.value = f"Bulk Upload Complete! {success_count} songs added."
                status_in.color = "green"
                page.update()
                # Refresh local cache in background
                cloud_sync()
                # Instead of closing, refresh the management list if visible
                time.sleep(2)
                show_admin_editor()

        threading.Thread(target=background_processing, daemon=True).start()


    # Use a hidden container for picker if needed, but overlay is better
    # We will append it to overlay only when the Admin dashboard is opened to avoid the red screen on home


    try:
        seed_default_songs()
    except Exception as e:
        page.add(ft.SafeArea(ft.Text(f"Initialization Error: {e}", color="red")))
        return

    # ---- State ----
    current_tab = [0]  # 0=Tamil, 1=Telugu
    search_query = [""]
    search_open = [False]
    is_syncing = [False]

    def refresh_ui():
        try:
            show_home()
        except:
            pass

    def start_auto_sync():
        if is_syncing[0]:
            return
        is_syncing[0] = True
        try:
            c = cloud_sync()
            if c > 0:
                show_home()
        finally:
            is_syncing[0] = False

    # ---- Build song list ----
    def make_song_list(lang, query=""):
        songs = get_filtered(lang, query)
        if not songs:
            return ft.Container(
                content=ft.Column(
                    controls=[
                        ft.Icon(ft.Icons.CLOUD_OFF, size=50, color="#B0BEC5"),
                        ft.Text(
                            "No songs found.\nChecking for updates...",
                            text_align=ft.TextAlign.CENTER,
                            color="#78909C",
                            size=16,
                        ),
                    ],
                    horizontal_alignment=ft.CrossAxisAlignment.CENTER,
                ),
                alignment=ft.alignment.Alignment(0, 0),
                padding=40,
            )
        
        col = ft.Column(spacing=10, scroll=ft.ScrollMode.AUTO, expand=True)
        for s in songs:
            fav = s.get("is_favorite", False)
            col.controls.append(
                ft.Container(
                    content=ft.Row(
                        controls=[
                            ft.Container(
                                content=ft.Icon(ft.Icons.MUSIC_NOTE, color="#3F51B5", size=22),
                                padding=10,
                                bgcolor="#E8EAF6",
                                border_radius=10,
                            ),
                            ft.Column(
                                controls=[
                                    ft.Text(s["title"], weight=ft.FontWeight.BOLD, size=16, color="#263238"),
                                    ft.Text(s["language"].capitalize(), size=12, color="#90A4AE"),
                                ],
                                expand=True,
                                spacing=0,
                            ),
                            ft.Icon(ft.Icons.FAVORITE, color="#FF5252", size=22) if fav else ft.Icon(ft.Icons.CHEVRON_RIGHT, color="#CFD8DC", size=24),
                        ],
                        vertical_alignment=ft.CrossAxisAlignment.CENTER,
                    ),
                    padding=ft.Padding(left=16, right=16, top=12, bottom=12),
                    border_radius=15,
                    bgcolor="white",
                    on_click=lambda e, sid=s["id"]: show_song(sid),
                    shadow=ft.BoxShadow(
                        blur_radius=10,
                        color=ft.Colors.with_opacity(0.05, "black"),
                        offset=ft.Offset(0, 4)
                    ),
                    ink=True,
                )
            )
        return col

    # ---- Tab Bar ----
    def make_tab_bar():
        def select_tab(idx):
            current_tab[0] = idx
            show_home()

        tamil_active = current_tab[0] == 0
        telugu_active = current_tab[0] == 1

        return ft.Container(
             padding=ft.Padding(left=10, right=10, top=5, bottom=5),
             bgcolor="#E1E5F2",
             border_radius=25,
             content=ft.Row(
                controls=[
                    ft.Container(
                        content=ft.Text("Tamil", size=14, weight=ft.FontWeight.BOLD if tamil_active else ft.FontWeight.NORMAL, color="white" if tamil_active else "#5C6BC0"),
                        bgcolor="#3F51B5" if tamil_active else "transparent",
                        border_radius=20,
                        padding=ft.Padding(left=20, right=20, top=8, bottom=8),
                        on_click=lambda e: select_tab(0),
                        expand=True,
                        alignment=ft.alignment.Alignment(0, 0),
                    ),
                    ft.Container(
                        content=ft.Text("Telugu", size=14, weight=ft.FontWeight.BOLD if telugu_active else ft.FontWeight.NORMAL, color="white" if telugu_active else "#5C6BC0"),
                        bgcolor="#3F51B5" if telugu_active else "transparent",
                        border_radius=20,
                        padding=ft.Padding(left=20, right=20, top=8, bottom=8),
                        on_click=lambda e: select_tab(1),
                        expand=True,
                        alignment=ft.alignment.Alignment(0, 0),
                    ),
                ],
                spacing=5,
            )
        )

    # ---- HOME SCREEN ----
    def show_home():
        try:
            lang = "tamil" if current_tab[0] == 0 else "telugu"
            song_list = make_song_list(lang, search_query[0])

            header = ft.Container(
                content=ft.Row(
                    controls=[
                        ft.Row([
                            ft.Image(src="icon.png", width=40, height=40, border_radius=8),
                            ft.Column([
                                ft.Text("GGGM", size=24, weight=ft.FontWeight.BOLD, color="white"),
                                ft.Text("Songs of Worship", size=12, color="#C5CAE9"),
                            ], spacing=2),
                        ], spacing=10),
                        ft.Row(
                            controls=[
                                ft.IconButton(ft.Icons.SEARCH, icon_color="white", on_click=lambda e: toggle_search(e)),
                                ft.IconButton(ft.Icons.SETTINGS, icon_color="white", on_click=lambda e: show_admin()),
                            ],
                            spacing=0,
                        ),
                    ],
                    alignment=ft.MainAxisAlignment.SPACE_BETWEEN,
                ),
                gradient=ft.LinearGradient(
                    begin=ft.alignment.Alignment(-1, -1),
                    end=ft.alignment.Alignment(1, 1),
                    colors=["#3F51B5", "#7986CB"]
                ),
                padding=ft.Padding(left=20, right=10, top=20, bottom=20),
            )

            controls = [header]

            if search_open[0]:
                controls.append(
                    ft.Container(
                        content=ft.TextField(
                            hint_text="Search song title...",
                            on_change=on_search_change,
                            value=search_query[0],
                            autofocus=True,
                            border_color="#C5CAE9",
                            border_radius=15,
                            prefix_icon=ft.Icons.SEARCH,
                            text_size=14,
                            bgcolor="white",
                        ),
                        padding=ft.Padding(left=16, right=16, top=10, bottom=10),
                    )
                )

            controls.append(
                ft.Container(
                    content=make_tab_bar(),
                    padding=ft.Padding(left=20, right=20, top=15, bottom=10),
                )
            )

            controls.append(
                ft.Container(content=song_list, padding=15, expand=True)
            )

            page.controls.clear()
            page.add(ft.SafeArea(content=ft.Column(controls=controls, spacing=0, expand=True), expand=True))
            page.update()
        except Exception as e:
            page.controls.clear()
            page.add(ft.SafeArea(ft.Text(f"Home error: {e}", color="red")))
            page.update()

    def on_search_change(e):
        search_query[0] = e.control.value
        show_home()

    def toggle_search(e):
        search_open[0] = not search_open[0]
        if not search_open[0]:
            search_query[0] = ""
        show_home()

    # ---- SONG DETAIL SCREEN ----
    def show_song(song_id):
        try:
            song = next((s for s in SONGS if s["id"] == song_id), None)
            if not song: return

            font_size = [18]
            lyrics_text = song.get("lyrics", "")
            if not lyrics_text:
                lyrics_text = "[No lyrics content found in database]"
            
            lyrics_widget = ft.Text(lyrics_text, size=font_size[0], selectable=True, color="#37474F")

            def update_font(delta):
                font_size[0] = max(10, min(40, font_size[0] + delta))
                lyrics_widget.size = font_size[0]
                page.update()

            def toggle_fav(e):
                song["is_favorite"] = not song.get("is_favorite", False)
                show_song(song_id)

            is_fav = song.get("is_favorite", False)

            page.controls.clear()
            page.add(
                ft.SafeArea(
                    content=ft.Column(
                        controls=[
                            ft.Container(
                                content=ft.Row(
                                    controls=[
                                        ft.IconButton(ft.Icons.ARROW_BACK, icon_color="white", on_click=lambda e: show_home()),
                                        ft.Text(song["title"], size=16, color="white", weight=ft.FontWeight.BOLD, expand=True, max_lines=1, overflow=ft.TextOverflow.ELLIPSIS),
                                        ft.Row([
                                            ft.IconButton(ft.Icons.REMOVE_CIRCLE_OUTLINE, icon_color="white", on_click=lambda e: update_font(-2)),
                                            ft.IconButton(ft.Icons.ADD_CIRCLE_OUTLINE, icon_color="white", on_click=lambda e: update_font(2)),
                                            ft.IconButton(ft.Icons.FAVORITE if is_fav else ft.Icons.FAVORITE_BORDER, icon_color="#FFCDD2" if is_fav else "white", on_click=toggle_fav),
                                        ], spacing=0)
                                    ],
                                ),
                                bgcolor="#3F51B5",
                                padding=ft.Padding(left=5, right=5, top=5, bottom=5),
                            ),
                            ft.Container(
                                content=ft.Column(
                                    controls=[
                                        ft.Text(song["title"], size=22, weight=ft.FontWeight.BOLD, color="#1A237E"),
                                        ft.Divider(color="#E8EAF6", height=30),
                                        lyrics_widget,
                                        ft.Container(height=50),
                                    ],
                                    scroll=ft.ScrollMode.AUTO,
                                ),
                                padding=25,
                                bgcolor="white",
                                expand=True,
                            ),
                        ],
                        spacing=0,
                        expand=True,
                    ),
                    expand=True,
                )
            )
            page.update()
        except Exception as e:
            page.controls.clear()
            page.add(ft.SafeArea(ft.Text(f"Song detail error: {e}", color="red")))
            page.update()

    # ---- ADMIN / SYNC SCREEN ----
    def show_admin():
        try:
            status = ft.Text("Ready to sync.", color="#43A047", size=14)

            def do_sync(e):
                status.value = "Syncing from cloud..."
                status.color = "#3F51B5"
                page.update()
                try:
                    c = cloud_sync()
                    status.value = f"Success! {c} songs updated."
                    status.color = "#43A047"
                except Exception as ex:
                    status.value = f"Sync failed: {ex}"
                    status.color = "#E53935"
                page.update()

            page.controls.clear()
            page.add(
                ft.SafeArea(
                    content=ft.Column(
                        controls=[
                            ft.Container(
                                content=ft.Row([
                                    ft.IconButton(ft.Icons.ARROW_BACK, icon_color="white", on_click=lambda e: show_home()),
                                    ft.Text("Settings", size=20, color="white", weight=ft.FontWeight.BOLD),
                                ]),
                                bgcolor="#3F51B5",
                                padding=ft.Padding(left=5, right=5, top=10, bottom=10),
                            ),
                            ft.Container(
                                content=ft.Column(
                                    controls=[
                                        ft.Text("App Appearance", size=24, weight=ft.FontWeight.BOLD, color="#1A237E"),
                                        ft.Text("Current App Logo", size=14, color="#546E7A"),
                                        ft.Container(height=10),
                                        ft.Container(
                                            content=ft.Image(src="icon.png", width=100, height=100, border_radius=15),
                                            bgcolor="white",
                                            padding=10,
                                            border_radius=20,
                                            shadow=ft.BoxShadow(blur_radius=10, color=ft.Colors.with_opacity(0.1, "black"))
                                        ),
                                        ft.Container(height=10),
                                        ft.ElevatedButton(
                                            content=ft.Row([ft.Icon(ft.Icons.LOCK), ft.Text("OPEN MASTER ADMIN PANEL")], alignment=ft.MainAxisAlignment.CENTER),
                                            on_click=lambda e: show_login(),
                                            height=50,
                                            style=ft.ButtonStyle(
                                                color="#3F51B5",
                                                bgcolor="#E1E5F2",
                                                shape=ft.RoundedRectangleBorder(radius=15),
                                            )
                                        ),
                                        ft.Divider(height=40, color="#E8EAF6"),
                                        ft.Text("Cloud Sync", size=24, weight=ft.FontWeight.BOLD, color="#1A237E"),
                                        ft.Text("Update your song library with the latest lyrics from the server.", size=14, color="#546E7A"),
                                        ft.Container(height=20),
                                        ft.ElevatedButton(
                                            content=ft.Row([ft.Icon(ft.Icons.SYNC), ft.Text("SYNC NOW")], alignment=ft.MainAxisAlignment.CENTER),
                                            on_click=do_sync,
                                            height=50,
                                            style=ft.ButtonStyle(
                                                color="white",
                                                bgcolor="#3F51B5",
                                                shape=ft.RoundedRectangleBorder(radius=15),
                                            )
                                        ),
                                        ft.Container(height=15),
                                        status,
                                    ],
                                    scroll=ft.ScrollMode.AUTO,
                                ),
                                padding=30,
                                expand=True,
                            ),
                        ],
                        spacing=0,
                        expand=True,
                    ),
                    expand=True,
                )
            )
            page.update()
        except Exception as e:
            page.controls.clear()
            page.add(ft.SafeArea(ft.Text(f"Settings error: {e}", color="red")))
            page.update()

    # ---- ADMIN LOGIN SCREEN ----
    def show_login():
        user_field = ft.TextField(label="Username", border_radius=15)
        pass_field = ft.TextField(label="Password", password=True, can_reveal_password=True, border_radius=15)
        error_text = ft.Text("", color="red")

        def attempt_login(e):
            if user_field.value == "admin" and pass_field.value == "Grace@2024":
                show_admin_editor()
            else:
                error_text.value = "Invalid credentials. Please try again."
                page.update()

        page.controls.clear()
        page.add(
            ft.SafeArea(
                content=ft.Column(
                    controls=[
                        ft.Container(
                            content=ft.Row([
                                ft.IconButton(ft.Icons.ARROW_BACK, icon_color="white", on_click=lambda e: show_admin()),
                                ft.Text("Admin Authentication", size=20, color="white", weight=ft.FontWeight.BOLD),
                            ]),
                            bgcolor="#3F51B5",
                            padding=ft.Padding(left=5, right=5, top=10, bottom=10),
                        ),
                        ft.Container(
                            content=ft.Column(
                                controls=[
                                    ft.Icon(ft.Icons.LOCK_PERSON, size=80, color="#3F51B5"),
                                    ft.Text("Master Access", size=28, weight=ft.FontWeight.BOLD, color="#1A237E"),
                                    ft.Text("Login to add or manage song lyrics.", color="#546E7A"),
                                    ft.Container(height=20),
                                    user_field,
                                    pass_field,
                                    error_text,
                                    ft.Container(height=10),
                                    ft.ElevatedButton(
                                        "Login", on_click=attempt_login, width=float("inf"), height=50,
                                        style=ft.ButtonStyle(bgcolor="#3F51B5", color="white", shape=ft.RoundedRectangleBorder(radius=15))
                                    ),
                                ],
                                horizontal_alignment=ft.CrossAxisAlignment.CENTER,
                            ),
                            padding=40,
                        )
                    ],
                    spacing=0,
                    expand=True,
                ),
                expand=True,
            )
        )
        page.update()

    # ---- ADMIN EDITOR (DASHBOARD) SCREEN ----
    def show_admin_editor():
        title_in = ft.TextField(label="Song Title", border_radius=12)
        lang_in = ft.Dropdown(
            label="Language",
            options=[ft.dropdown.Option("tamil"), ft.dropdown.Option("telugu")],
            border_radius=12,
            value="tamil"
        )
        num_in = ft.TextField(label="Song Number (Optional)", border_radius=12)
        lyrics_in = ft.TextField(label="Full Lyrics", multiline=True, min_lines=8, max_lines=15, border_radius=12)
        status_in = ft.Text("", color="blue")
        
        # Save references for bulk upload
        page.admin_status = status_in
        page.admin_lang = lang_in.value

        def push_to_cloud(e):
            if not title_in.value or not lyrics_in.value:
                status_in.value = "Error: Title and Lyrics are required."
                status_in.color = "red"
                page.update()
                return

            status_in.value = "Pushing to database..."
            status_in.color = "blue"
            page.update()

            try:
                import uuid
                song_id = f"{lang_in.value[:2]}_{uuid.uuid4().hex[:8]}"
                document = {
                    "fields": {
                        "title": {"stringValue": title_in.value},
                        "language": {"stringValue": lang_in.value},
                        "number": {"stringValue": num_in.value or ""},
                        "lyrics": {"stringValue": lyrics_in.value.replace("\n", "\\n")},
                        "category": {"stringValue": "General"},
                        "composer": {"stringValue": "Unknown"}
                    }
                }
                
                req = urllib.request.Request(
                    f"{FIRESTORE_URL}/{song_id}",
                    data=json.dumps(document).encode("utf-8"),
                    headers={"Content-Type": "application/json"},
                    method="PATCH"
                )
                with urllib.request.urlopen(req) as r:
                    if r.status == 200:
                        status_in.value = f"Success! Uploaded as {song_id}."
                        status_in.color = "green"
                        title_in.value = ""
                        lyrics_in.value = ""
                        num_in.value = ""
                    else:
                        status_in.value = f"Failed (HTTP {r.status})"
                        status_in.color = "red"
            except Exception as ex:
                status_in.value = f"Error: {ex}"
                status_in.color = "red"
            page.update()

        page.controls.clear()
        page.add(
            ft.SafeArea(
                content=ft.Column(
                    controls=[
                        ft.Container(
                            content=ft.Row([
                                ft.IconButton(ft.Icons.ARROW_BACK, icon_color="white", on_click=lambda e: show_login()),
                                ft.Text("Cloud Management", size=20, color="white", weight=ft.FontWeight.BOLD),
                            ]),
                            bgcolor="#3F51B5",
                            padding=ft.Padding(left=5, right=5, top=10, bottom=10),
                        ),
                        ft.Container(
                            content=ft.Column(
                                controls=[
                                    ft.Text("Manage / Delete Songs", size=24, weight=ft.FontWeight.BOLD, color="#1A237E"),
                                    ft.Text("Deletions are permanent and sync to all users.", size=12, color="#546E7A"),
                                    ft.Text(f"Total Songs: {len(SONGS)}", size=14, color="#546E7A"),
                                ] + [
                                    ft.Container(
                                        content=ft.Row([
                                            ft.Column([
                                                ft.Text(s["title"], weight=ft.FontWeight.BOLD, size=16),
                                                ft.Text(f"ID: {s['id']} | Lang: {s['language']}", size=12, color="#78909C"),
                                            ], expand=True),
                                            ft.IconButton(
                                                ft.Icons.DELETE_OUTLINE, 
                                                icon_color="red", 
                                                tooltip="Delete from Cloud",
                                                on_click=lambda e, sid=s['id']: delete_song(sid)
                                            ),
                                        ]),
                                        padding=15,
                                        bgcolor="#F8F9FB",
                                        border_radius=12,
                                        border=ft.border.all(1, "#E8EAF6"),
                                    ) for s in sorted(SONGS, key=lambda x: x["title"])[:50] # Show top 50
                                ] + [
                                    ft.Text("Showing first 50 songs. Use search for more.", size=12, italic=True) if len(SONGS) > 50 else ft.Container()
                                ],
                                scroll=ft.ScrollMode.AUTO,
                                spacing=15,
                            ),
                            padding=20,
                            expand=True,
                        )
                    ],
                    spacing=0,
                    expand=True,
                ),
                expand=True,
            )
        )
        page.update()

    def delete_song(sid):
        def do_delete():
            try:
                # Use our new Cloudflare /delete endpoint
                data = json.dumps({"id": sid}).encode("utf-8")
                req = urllib.request.Request(f"{FIRESTORE_URL}/delete", data=data, method="POST")
                with urllib.request.urlopen(req) as r:
                    if r.status == 200:
                        cloud_sync() # Refresh locally
                        show_admin_editor() # Refresh UI
            except Exception as ex:
                print(f"Delete error: {ex}")
        
        threading.Thread(target=do_delete, daemon=True).start()

    # ---- START ----

    show_home()
    # Trigger auto-sync in background
    threading.Thread(target=start_auto_sync, daemon=True).start()

ft.app(main, assets_dir="assets")
